import re
from pathlib import Path
from typing import Any

from .config import get_settings


def _clean(text: str) -> str:
    return re.sub(r"[`*_>#]", "", text).strip()


def _bullets(markdown: str, limit: int = 5) -> list[str]:
    items = []
    for line in markdown.splitlines():
        line = line.strip()
        if line.startswith(("-", "*")):
            value = _clean(line[1:])
            if value:
                items.append(value)
        if len(items) >= limit:
            break
    return items


def build_pptx(resource: dict[str, Any]) -> Path:
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.util import Inches, Pt
    except ImportError as error:
        raise RuntimeError("未安装 python-pptx，请执行 pip install -r backend/requirements.txt") from error

    content = resource["content"]
    profile = content.get("profile_snapshot", {})
    topic = resource["topic"]
    lecture = str(content.get("lecture", ""))
    path = str(content.get("learning_path", ""))
    code = str(content.get("code_case", ""))
    quiz = content.get("quiz", [])
    strategy = profile.get("learning_style", "个性化学习")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def style(slide) -> None:
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(248, 250, 252)
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.name = "Microsoft YaHei"
                    run.font.color.rgb = RGBColor(30, 41, 59)

    def add_title(title: str, subtitle: str) -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = title
        slide.placeholders[1].text = subtitle
        style(slide)
        slide.shapes.title.text_frame.paragraphs[0].runs[0].font.size = Pt(34)

    def add_content(title: str, items: list[str], note: str = "") -> None:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title
        body = slide.placeholders[1].text_frame
        body.clear()
        for index, item in enumerate(items or ["内容由资源包生成"]):
            paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
            paragraph.text = item
            paragraph.level = 0
            paragraph.font.size = Pt(21)
        if note:
            box = slide.shapes.add_textbox(Inches(0.8), Inches(6.55), Inches(11.8), Inches(0.4))
            box.text_frame.text = note
            box.text_frame.paragraphs[0].font.size = Pt(11)
            box.text_frame.paragraphs[0].font.color.rgb = RGBColor(109, 40, 217)
        style(slide)

    add_title(
        topic,
        f"{profile.get('display_name', '当前学习者')} · {strategy}",
    )
    add_content(
        "学习目标与个性化路径",
        _bullets(path) or [
            f"当前基础：{profile.get('knowledge_level', '未填写')}",
            f"重点补强：{'、'.join(profile.get('weak_points', []))}",
            f"资源偏好：{'、'.join(profile.get('resource_preference', []))}",
        ],
        "本页依据当前学生画像与历史薄弱点生成",
    )
    add_content("知识点讲解", _bullets(lecture, 6), "内容需结合知识库引用核验")
    code_lines = [
        _clean(line)
        for line in code.splitlines()
        if line.strip() and not line.strip().startswith(("#", "```"))
    ][:8]
    add_content("示例代码与实操任务", code_lines, "代码未在服务器执行，运行结果需人工核验")
    add_content(
        "针对性练习",
        [f"{index + 1}. {item.get('question')}" for index, item in enumerate(quiz[:4])],
        "选择/判断精确判分；简答与代码题为 MVP 评分",
    )
    if "考试" in strategy:
        summary = ["复盘高频考点与易错 API", "完成限时练习", "整理一页错题清单"]
    elif "项目" in strategy:
        summary = ["用真实 CSV 替换示例数据", "输出清洗前后对比", "提交可复现分析报告"]
    else:
        summary = ["复述数据清洗六步流程", "修改最小代码示例", "24 小时后完成一次复测"]
    add_content("总结与下一步建议", summary, "测评结果将写回画像并更新 Dashboard")

    target = get_settings().exports_dir / f"resource-{resource['id']}.pptx"
    prs.save(target)
    return target
